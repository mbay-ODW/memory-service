"""Text extraction for the document-upload feature -- deliberately extraction-only, no
retention of the uploaded bytes anywhere (not in git, not on disk after this call returns).
Murat already runs Paperless-ngx for OCR'd document storage/retention; this exists purely to
save re-typing a digitally-born document's content into a new entry, not to become a second
document store. None of the libraries here do OCR -- a scanned/image-only PDF extracts to
empty or garbage text either way, which is fine: that's Paperless's job, not this one's."""

import asyncio
import csv
import io
from html.parser import HTMLParser

ALLOWED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx", ".txt", ".md", ".html", ".htm", ".csv"}


class ExtractionError(Exception):
    pass


class _TextOnlyHTMLParser(HTMLParser):
    _SKIP_TAGS = {"script", "style"}

    def __init__(self):
        super().__init__()
        self._chunks: list[str] = []
        self._skip_depth = 0

    def handle_starttag(self, tag: str, attrs) -> None:
        if tag in self._SKIP_TAGS:
            self._skip_depth += 1

    def handle_endtag(self, tag: str) -> None:
        if tag in self._SKIP_TAGS and self._skip_depth > 0:
            self._skip_depth -= 1

    def handle_data(self, data: str) -> None:
        if self._skip_depth:
            return
        stripped = data.strip()
        if stripped:
            self._chunks.append(stripped)

    def get_text(self) -> str:
        return "\n".join(self._chunks)


def _extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    return "\n\n".join(page.extract_text() or "" for page in reader.pages)


def _extract_docx(data: bytes) -> str:
    import docx

    document = docx.Document(io.BytesIO(data))
    return "\n\n".join(p.text for p in document.paragraphs if p.text.strip())


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(data))
    chunks = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                chunks.append(shape.text_frame.text.strip())
    return "\n\n".join(chunks)


def _extract_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    chunks = []
    for sheet in workbook.worksheets:
        chunks.append(f"## {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                chunks.append(" | ".join(cells))
    return "\n".join(chunks)


def _extract_csv(data: bytes) -> str:
    text = data.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    return "\n".join(" | ".join(row) for row in reader if row)


def _extract_html(data: bytes) -> str:
    parser = _TextOnlyHTMLParser()
    parser.feed(data.decode("utf-8", errors="replace"))
    return parser.get_text()


def _extract_plain(data: bytes) -> str:
    return data.decode("utf-8", errors="replace")


_EXTRACTORS = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".xlsx": _extract_xlsx,
    ".csv": _extract_csv,
    ".html": _extract_html,
    ".htm": _extract_html,
    ".txt": _extract_plain,
    ".md": _extract_plain,
}


def _extract_sync(filename: str, data: bytes) -> str:
    suffix = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    extractor = _EXTRACTORS.get(suffix)
    if extractor is None:
        raise ExtractionError(
            f"Dateityp {suffix or '(ohne Endung)'!r} wird nicht unterstützt. "
            f"Erlaubt: {', '.join(sorted(ALLOWED_EXTENSIONS))}"
        )
    try:
        text = extractor(data)
    except ExtractionError:
        raise
    except Exception as exc:
        raise ExtractionError(f"Datei konnte nicht gelesen werden: {exc}") from exc
    if not text.strip():
        raise ExtractionError(
            "Es konnte kein Text aus der Datei extrahiert werden (z. B. weil es sich um ein "
            "gescanntes/bildbasiertes Dokument ohne eingebetteten Text handelt)."
        )
    return text


async def extract_text(filename: str, data: bytes) -> str:
    return await asyncio.to_thread(_extract_sync, filename, data)
