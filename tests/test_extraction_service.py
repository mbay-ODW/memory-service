import io

import pytest

from app.services import extraction as extraction_service


def _build_minimal_pdf(text: str) -> bytes:
    """Hand-built minimal single-page PDF with a real xref table -- pypdf needs a correct
    startxref/xref, so this can't just be free-form PDF syntax; offsets are computed
    programmatically to avoid manual arithmetic errors."""
    objects = [
        b"<</Type/Catalog/Pages 2 0 R>>",
        b"<</Type/Pages/Kids[3 0 R]/Count 1>>",
        b"<</Type/Page/Parent 2 0 R/Resources<</Font<</F1 4 0 R>>>>/MediaBox[0 0 200 200]/Contents 5 0 R>>",
        b"<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>",
    ]
    stream = f"BT /F1 24 Tf 10 100 Td ({text}) Tj ET".encode()
    objects.append(b"<</Length " + str(len(stream)).encode() + b">>\nstream\n" + stream + b"\nendstream")

    out = io.BytesIO()
    out.write(b"%PDF-1.4\n")
    offsets = [0]
    for i, obj in enumerate(objects, start=1):
        offsets.append(out.tell())
        out.write(f"{i} 0 obj".encode() + obj + b"endobj\n")
    xref_pos = out.tell()
    out.write(f"xref\n0 {len(objects) + 1}\n".encode())
    out.write(b"0000000000 65535 f \n")
    for off in offsets[1:]:
        out.write(f"{off:010d} 00000 n \n".encode())
    out.write(f"trailer<</Size {len(objects) + 1}/Root 1 0 R>>\nstartxref\n{xref_pos}\n%%EOF".encode())
    return out.getvalue()


async def test_extract_pdf():
    text = await extraction_service.extract_text("test.pdf", _build_minimal_pdf("Hallo PDF Test"))
    assert "Hallo PDF Test" in text


async def test_extract_docx():
    import docx

    document = docx.Document()
    document.add_paragraph("Hallo DOCX Test")
    buf = io.BytesIO()
    document.save(buf)

    text = await extraction_service.extract_text("test.docx", buf.getvalue())
    assert "Hallo DOCX Test" in text


async def test_extract_pptx():
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    textbox.text_frame.text = "Hallo PPTX Test"
    buf = io.BytesIO()
    presentation.save(buf)

    text = await extraction_service.extract_text("test.pptx", buf.getvalue())
    assert "Hallo PPTX Test" in text


async def test_extract_xlsx():
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet["A1"] = "Hallo XLSX Test"
    buf = io.BytesIO()
    workbook.save(buf)

    text = await extraction_service.extract_text("test.xlsx", buf.getvalue())
    assert "Hallo XLSX Test" in text


async def test_extract_csv():
    data = "Name,Wert\nHallo CSV Test,42\n".encode("utf-8")
    text = await extraction_service.extract_text("test.csv", data)
    assert "Hallo CSV Test" in text


async def test_extract_html():
    data = b"<html><body><h1>Hallo HTML Test</h1><script>ignored();</script></body></html>"
    text = await extraction_service.extract_text("test.html", data)
    assert "Hallo HTML Test" in text
    assert "ignored" not in text


async def test_extract_plain_text():
    text = await extraction_service.extract_text("test.txt", "Hallo TXT Test".encode("utf-8"))
    assert text == "Hallo TXT Test"


async def test_extract_rejects_unsupported_extension():
    with pytest.raises(extraction_service.ExtractionError):
        await extraction_service.extract_text("test.exe", b"whatever")


async def test_extract_raises_on_empty_extracted_text():
    with pytest.raises(extraction_service.ExtractionError):
        await extraction_service.extract_text("empty.txt", b"   \n\n  ")
